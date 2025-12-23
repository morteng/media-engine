"""
PowerPoint Builder

Generates PPTX presentations from Markdown or YAML definitions.
Applies theme styling automatically.

Features:
- Title slides
- Content slides with bullets
- Two-column layouts
- Image slides
- Quote slides
- Section dividers
"""

import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import TYPE_CHECKING, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

if TYPE_CHECKING:
    from ..brand import BrandContext
    from ..core.theme import Theme


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGBColor."""
    hex_color = hex_color.lstrip("#")
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


@dataclass
class SlideContent:
    """Content for a single slide."""

    type: str  # title, content, two_column, image, quote, section
    title: str = ""
    subtitle: str = ""
    bullets: List[str] = field(default_factory=list)
    left_bullets: List[str] = field(default_factory=list)
    right_bullets: List[str] = field(default_factory=list)
    image_path: Optional[Path] = None
    image_caption: str = ""
    quote: str = ""
    quote_author: str = ""
    notes: str = ""


class PPTXBuilder:
    """Builds PowerPoint presentations with theme styling."""

    def __init__(
        self,
        theme: "Theme" = None,
        logo_path: Optional[Path] = None,
        brand: "BrandContext" = None,
    ):
        """
        Initialize PPTX builder.

        Args:
            theme: Legacy Theme for styling (deprecated, use brand instead)
            logo_path: Legacy logo path (deprecated, use brand instead)
            brand: BrandContext for unified brand access (recommended)
        """
        self.brand = brand
        self._logo_path = logo_path

        if brand:
            # Use BrandContext (recommended)
            self.color_primary = hex_to_rgb(brand.get_brand_color("primary"))
            self.color_secondary = hex_to_rgb(brand.get_brand_color("secondary"))
            self.color_accent = hex_to_rgb(brand.get_brand_color("accent"))
            self.color_bg = hex_to_rgb(brand.get_color("background.primary") or "#ffffff")
            self.color_muted = hex_to_rgb(brand.get_color("text.muted") or "#666666")

            self.font_heading = brand.get_system_font("heading")
            self.font_body = brand.get_system_font("body")
        else:
            # Legacy Theme support
            from ..core.theme import COPPER_AND_CREAM

            self.theme = theme or COPPER_AND_CREAM

            self.color_primary = hex_to_rgb(self.theme.colors.primary)
            self.color_secondary = hex_to_rgb(self.theme.colors.secondary)
            self.color_accent = hex_to_rgb(self.theme.colors.accent)
            self.color_bg = hex_to_rgb(self.theme.colors.background)
            self.color_muted = hex_to_rgb(self.theme.colors.muted)

            self.font_heading = self._get_system_font(self.theme.typography.heading)
            self.font_body = self._get_system_font(self.theme.typography.body)

        # Slide dimensions (16:9)
        self.slide_width = Inches(13.333)
        self.slide_height = Inches(7.5)

        # Create presentation
        self.prs = Presentation()
        self.prs.slide_width = self.slide_width
        self.prs.slide_height = self.slide_height

    @property
    def logo_path(self) -> Optional[Path]:
        """Get logo path, preferring BrandContext if available."""
        if self.brand:
            # Use brand context to get PNG (auto-converts from SVG)
            return self.brand.get_logo("primary", "png", size=200)
        return self._logo_path

    def _get_logo_for_pptx(self) -> Optional[Path]:
        """Get logo in a format suitable for PPTX (PNG/JPG, not SVG)."""
        # If using BrandContext, it handles conversion
        if self.brand:
            return self.brand.get_logo("primary", "png", size=200)

        # Legacy path handling
        if not self._logo_path or not self._logo_path.exists():
            return None

        # If it's already PNG or JPG, use directly
        if self._logo_path.suffix.lower() in [".png", ".jpg", ".jpeg"]:
            return self._logo_path

        # For SVG, try to convert to PNG
        if self._logo_path.suffix.lower() == ".svg":
            try:
                import tempfile

                import cairosvg

                # Create temp PNG file
                png_path = Path(tempfile.gettempdir()) / f"{self._logo_path.stem}_logo.png"
                if not png_path.exists():
                    cairosvg.svg2png(
                        url=str(self._logo_path),
                        write_to=str(png_path),
                        output_height=200,  # High resolution for quality
                    )
                return png_path
            except ImportError:
                return None  # cairosvg not available
            except Exception:
                return None

        return None

    def _add_logo(self, slide, position: str = "top-left"):
        """Add logo to slide if logo_path is set."""
        logo_file = self._get_logo_for_pptx()
        if not logo_file:
            return

        # Logo positioning
        if position == "top-left":
            left, top = Inches(0.5), Inches(0.3)
        elif position == "top-right":
            left, top = Inches(11.5), Inches(0.3)
        elif position == "bottom-right":
            left, top = Inches(11.5), Inches(6.8)
        else:
            left, top = Inches(0.5), Inches(0.3)

        try:
            slide.shapes.add_picture(str(logo_file), left, top, height=Inches(0.6))
        except Exception:
            pass  # Skip logo if there's an issue

    def _get_system_font(self, font_name: str) -> str:
        """Map theme font to system font for compatibility."""
        font_map = {
            "Fraunces": "Georgia",
            "Source Sans 3": "Helvetica Neue",
            "Inter": "Helvetica Neue",
            "JetBrains Mono": "Menlo",
        }
        return font_map.get(font_name, font_name)

    def _set_background(self, slide, color: RGBColor):
        """Set slide background color."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_text_box(
        self,
        slide,
        text: str,
        left: float,
        top: float,
        width: float,
        height: float,
        font_name: str = None,
        font_size: int = 18,
        font_color: RGBColor = None,
        bold: bool = False,
        alignment: PP_ALIGN = PP_ALIGN.LEFT,
    ):
        """Add a text box to a slide."""
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text
        p.font.name = font_name or self.font_body
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color or self.color_primary
        p.font.bold = bold
        p.alignment = alignment

        return txBox

    def _add_bullet_list(
        self,
        slide,
        items: List[str],
        left: float,
        top: float,
        width: float,
        height: float,
        font_size: int = 18,
    ):
        """Add a bulleted list to a slide."""
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = f"• {item}"
            p.font.name = self.font_body
            p.font.size = Pt(font_size)
            p.font.color.rgb = self.color_primary
            p.level = 0
            p.space_before = Pt(8)
            p.space_after = Pt(4)

        return txBox

    def add_title_slide(self, title: str, subtitle: str = "", notes: str = ""):
        """Add a title slide."""
        blank_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(blank_layout)
        self._set_background(slide, self.color_bg)

        # Logo at top center for title slides
        logo_file = self._get_logo_for_pptx()
        if logo_file:
            try:
                slide.shapes.add_picture(
                    str(logo_file), Inches(5.5), Inches(1.0), height=Inches(1.2)
                )
            except Exception:
                pass

        # Title (positioned lower if logo exists)
        title_top = 2.8 if self.logo_path else 2.5
        self._add_text_box(
            slide,
            title,
            left=0.8,
            top=title_top,
            width=11.7,
            height=1.5,
            font_name=self.font_heading,
            font_size=54,
            font_color=self.color_primary,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        # Subtitle
        if subtitle:
            self._add_text_box(
                slide,
                subtitle,
                left=0.8,
                top=title_top + 1.5,
                width=11.7,
                height=0.8,
                font_name=self.font_body,
                font_size=24,
                font_color=self.color_muted,
                alignment=PP_ALIGN.CENTER,
            )

        # Accent bar
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(title_top + 1.3), Inches(2.3), Inches(0.05)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.color_accent
        shape.line.fill.background()

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        return slide

    def add_content_slide(
        self,
        title: str,
        bullets: List[str],
        notes: str = "",
    ):
        """Add a content slide with bullets."""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self._set_background(slide, self.color_bg)

        # Title
        self._add_text_box(
            slide,
            title,
            left=0.8,
            top=0.5,
            width=11.7,
            height=0.8,
            font_name=self.font_heading,
            font_size=36,
            font_color=self.color_primary,
            bold=True,
        )

        # Accent underline
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.03)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.color_accent
        shape.line.fill.background()

        # Bullets
        if bullets:
            self._add_bullet_list(
                slide,
                bullets,
                left=0.8,
                top=1.5,
                width=11.7,
                height=5.5,
                font_size=20,
            )

        # Small logo at bottom right
        self._add_logo(slide, "bottom-right")

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        return slide

    def add_two_column_slide(
        self,
        title: str,
        left_title: str,
        left_bullets: List[str],
        right_title: str,
        right_bullets: List[str],
        notes: str = "",
    ):
        """Add a two-column comparison slide."""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self._set_background(slide, self.color_bg)

        # Title
        self._add_text_box(
            slide,
            title,
            left=0.8,
            top=0.5,
            width=11.7,
            height=0.8,
            font_name=self.font_heading,
            font_size=36,
            font_color=self.color_primary,
            bold=True,
        )

        # Left column title
        self._add_text_box(
            slide,
            left_title,
            left=0.8,
            top=1.5,
            width=5.5,
            height=0.5,
            font_name=self.font_heading,
            font_size=24,
            font_color=self.color_accent,
            bold=True,
        )

        # Left column bullets
        self._add_bullet_list(
            slide,
            left_bullets,
            left=0.8,
            top=2.1,
            width=5.5,
            height=4.5,
            font_size=18,
        )

        # Right column title
        self._add_text_box(
            slide,
            right_title,
            left=7.0,
            top=1.5,
            width=5.5,
            height=0.5,
            font_name=self.font_heading,
            font_size=24,
            font_color=self.color_accent,
            bold=True,
        )

        # Right column bullets
        self._add_bullet_list(
            slide,
            right_bullets,
            left=7.0,
            top=2.1,
            width=5.5,
            height=4.5,
            font_size=18,
        )

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        return slide

    def add_image_slide(
        self,
        title: str,
        image_path: Path,
        caption: str = "",
        notes: str = "",
    ):
        """Add a slide with an image."""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self._set_background(slide, self.color_bg)

        # Title
        self._add_text_box(
            slide,
            title,
            left=0.8,
            top=0.5,
            width=11.7,
            height=0.8,
            font_name=self.font_heading,
            font_size=36,
            font_color=self.color_primary,
            bold=True,
        )

        # Image
        if image_path.exists():
            slide.shapes.add_picture(
                str(image_path), Inches(1.5), Inches(1.5), width=Inches(10.3), height=Inches(5.2)
            )

        # Caption
        if caption:
            self._add_text_box(
                slide,
                caption,
                left=0.8,
                top=6.9,
                width=11.7,
                height=0.4,
                font_size=14,
                font_color=self.color_muted,
                alignment=PP_ALIGN.CENTER,
            )

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        return slide

    def add_quote_slide(
        self,
        quote: str,
        author: str = "",
        notes: str = "",
    ):
        """Add a quote slide."""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self._set_background(slide, self.color_primary)

        # Quote
        self._add_text_box(
            slide,
            f'"{quote}"',
            left=1.5,
            top=2.0,
            width=10.3,
            height=3.0,
            font_name=self.font_heading,
            font_size=32,
            font_color=hex_to_rgb("#ffffff"),
            alignment=PP_ALIGN.CENTER,
        )

        # Author
        if author:
            self._add_text_box(
                slide,
                f"— {author}",
                left=1.5,
                top=5.0,
                width=10.3,
                height=0.5,
                font_size=20,
                font_color=self.color_accent,
                alignment=PP_ALIGN.CENTER,
            )

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        return slide

    def add_section_slide(
        self,
        title: str,
        subtitle: str = "",
        notes: str = "",
    ):
        """Add a section divider slide."""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self._set_background(slide, self.color_accent)

        # Section number/title
        self._add_text_box(
            slide,
            title,
            left=0.8,
            top=3.0,
            width=11.7,
            height=1.0,
            font_name=self.font_heading,
            font_size=48,
            font_color=hex_to_rgb("#ffffff"),
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        if subtitle:
            self._add_text_box(
                slide,
                subtitle,
                left=0.8,
                top=4.2,
                width=11.7,
                height=0.6,
                font_size=24,
                font_color=hex_to_rgb("#ffffff"),
                alignment=PP_ALIGN.CENTER,
            )

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        return slide

    def add_slide(self, content: SlideContent):
        """Add a slide based on SlideContent."""
        if content.type == "title":
            return self.add_title_slide(content.title, content.subtitle, content.notes)
        elif content.type == "content":
            return self.add_content_slide(content.title, content.bullets, content.notes)
        elif content.type == "two_column":
            return self.add_two_column_slide(
                content.title,
                "Left",
                content.left_bullets,
                "Right",
                content.right_bullets,
                content.notes,
            )
        elif content.type == "image":
            return self.add_image_slide(
                content.title, content.image_path, content.image_caption, content.notes
            )
        elif content.type == "quote":
            return self.add_quote_slide(content.quote, content.quote_author, content.notes)
        elif content.type == "section":
            return self.add_section_slide(content.title, content.subtitle, content.notes)
        else:
            return self.add_content_slide(content.title, content.bullets, content.notes)

    def build_from_markdown(self, markdown_path: Path) -> "PPTXBuilder":
        """
        Build presentation from Markdown file.

        Markdown format:
        ```
        # Title Slide
        Subtitle text

        ---

        ## Content Slide Title
        - Bullet 1
        - Bullet 2

        ---

        > Quote text
        > — Author

        ---
        ```
        """
        content = markdown_path.read_text()
        slides = content.split("\n---\n")

        for slide_content in slides:
            slide_content = slide_content.strip()
            if not slide_content:
                continue

            lines = slide_content.split("\n")

            # Title slide (# heading)
            if lines[0].startswith("# ") and not lines[0].startswith("## "):
                title = lines[0][2:].strip()
                subtitle = "\n".join(lines[1:]).strip() if len(lines) > 1 else ""
                self.add_title_slide(title, subtitle)

            # Quote (> blockquote)
            elif lines[0].startswith("> "):
                quote_lines = [line[2:] for line in lines if line.startswith("> ")]
                quote = " ".join(quote_lines)
                # Check for author
                if "—" in quote or "-" in quote:
                    parts = re.split(r"[—-]", quote, 1)
                    self.add_quote_slide(
                        parts[0].strip(), parts[1].strip() if len(parts) > 1 else ""
                    )
                else:
                    self.add_quote_slide(quote)

            # Content slide (## heading + bullets)
            elif lines[0].startswith("## "):
                title = lines[0][3:].strip()
                bullets = []
                for line in lines[1:]:
                    line = line.strip()
                    if line.startswith("- ") or line.startswith("* "):
                        bullets.append(line[2:])
                    elif line.startswith("• "):
                        bullets.append(line[2:])
                self.add_content_slide(title, bullets)

            # Section slide (### heading)
            elif lines[0].startswith("### "):
                title = lines[0][4:].strip()
                subtitle = "\n".join(lines[1:]).strip() if len(lines) > 1 else ""
                self.add_section_slide(title, subtitle)

        return self

    def build_from_yaml(self, yaml_path: Path) -> "PPTXBuilder":
        """
        Build presentation from YAML file.

        YAML format:
        ```yaml
        title: "Presentation Title"
        subtitle: "Optional subtitle"

        slides:
          - type: title
            title: "Title Slide"
            subtitle: "Subtitle"

          - type: content
            title: "Content Slide"
            bullets:
              - "Point 1"
              - "Point 2"

          - type: quote
            quote: "Quote text"
            author: "Author name"
        ```
        """
        import yaml

        with open(yaml_path, "r") as f:
            data = yaml.safe_load(f)

        # Title slide if defined at top level
        if "title" in data:
            self.add_title_slide(data["title"], data.get("subtitle", ""))

        # Process slides
        for slide_data in data.get("slides", []):
            slide_type = slide_data.get("type", "content")

            content = SlideContent(
                type=slide_type,
                title=slide_data.get("title", ""),
                subtitle=slide_data.get("subtitle", ""),
                bullets=slide_data.get("bullets", []),
                left_bullets=slide_data.get("left_bullets", []),
                right_bullets=slide_data.get("right_bullets", []),
                image_path=Path(slide_data["image"]) if "image" in slide_data else None,
                image_caption=slide_data.get("caption", ""),
                quote=slide_data.get("quote", ""),
                quote_author=slide_data.get("author", ""),
                notes=slide_data.get("notes", ""),
            )

            self.add_slide(content)

        return self

    def save(self, output_path: Path) -> Path:
        """Save presentation to file."""
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        return output_path


def build_pptx(
    source_path: Path,
    output_path: Path,
    theme: "Theme" = None,
) -> Path:
    """
    Build PPTX from Markdown or YAML source.

    Args:
        source_path: Path to .md or .yaml file
        output_path: Where to save .pptx
        theme: Optional theme for styling

    Returns:
        Path to generated PPTX
    """
    builder = PPTXBuilder(theme=theme)

    if source_path.suffix in [".yaml", ".yml"]:
        builder.build_from_yaml(source_path)
    else:
        builder.build_from_markdown(source_path)

    return builder.save(output_path)
